from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# --- FLAT ACADEMIC COLOR PALETTE ---
COLOR_BG = RGBColor(255, 255, 255)       # Pure White (Paper style)
COLOR_TEXT = RGBColor(40, 40, 40)        # Dark Grey (Softer than black)
COLOR_ACCENT = RGBColor(31, 78, 121)     # Flat Blue (Standard academic color)
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)
COLOR_RED = RGBColor(200, 50, 50)        # Flat Red for negatives
COLOR_GREEN = RGBColor(50, 160, 50)      # Flat Green for positives

# Helper to format text
def set_run_format(run, font_name="Arial", size=18, bold=False, color=COLOR_TEXT):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

# ==========================================
# SLIDE 1: Title & Team Members
# ==========================================
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)

# Force White Background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_BG

# Title
title = slide.shapes.title
title.text = "Term Project Proposal: Dungeon Game"
set_run_format(title.text_frame.paragraphs[0].runs[0], size=36, bold=True, color=COLOR_ACCENT)
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Subtitle
subtitle = slide.placeholders[1]
subtitle.text = "CSX3009 | LeetCode 174"
set_run_format(subtitle.text_frame.paragraphs[0].runs[0], size=20, color=COLOR_TEXT)
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Team Members (Add a flat text box)
left = Inches(1.5)
top = Inches(4.0)
width = Inches(7)
height = Inches(2.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Team Members"
p.font.bold = True
p.font.size = Pt(22)
p.font.color.rgb = COLOR_ACCENT
p.space_after = Pt(10)

# Team List
members = ["Puran Paodensakul 6611140", "ANUSON KHWANSAKUN 6610789"]
for member in members:
    p = tf.add_paragraph()
    p.text = member
    p.level = 1
    set_run_format(p.runs[0], size=18)

# ==========================================
# SLIDE 2: Problem Description
# ==========================================
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_BG

# Flat Header Bar
header_shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(1.0)) # Rectangle
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = COLOR_ACCENT
header_shape.line.fill.background()

# Header Text
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Problem Description"
set_run_format(p.runs[0], size=24, bold=True, color=RGBColor(255,255,255))

# Body Content
body_left = Inches(1.0)
body_top = Inches(1.5)
body_width = Inches(8.0)

txBox = slide.shapes.add_textbox(body_left, body_top, body_width, Inches(5))
tf = txBox.text_frame

points = [
    "The dungeon is represented as a 2D grid.",
    "Each cell contains positive or negative health values.",
    "The knight starts at the top-left cell.",
    "The knight can move only right or down.",
    "The goal is to reach the bottom-right cell.",
    "The knight dies if health becomes 0 or less."
]

for pt in points:
    p = tf.add_paragraph()
    p.text = pt
    p.level = 0
    set_run_format(p.runs[0], size=20)
    p.space_before = Pt(10)

# Objective Box
obj_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(8.0), Inches(1.5))
obj_box.fill.solid()
obj_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
tf = obj_box.text_frame
p = tf.add_paragraph()
p.text = "Objective: Find the minimum initial health required to guarantee survival."
set_run_format(p.runs[0], size=20, bold=True)

# ==========================================
# SLIDE 3: Input / Output Format
# ==========================================
slide = prs.slides.add_slide(slide_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_BG

# Flat Header
header_shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(1.0))
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = COLOR_ACCENT
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Input / Output Format"
set_run_format(p.runs[0], size=24, bold=True, color=RGBColor(255,255,255))

# Left Column: Description
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.0), Inches(5))
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = "Input:"
set_run_format(p.runs[0], size=22, bold=True)
p = tf.add_paragraph()
p.text = "A 2D integer array dungeon[m][n]"
set_run_format(p.runs[0], size=18)

p = tf.add_paragraph(); p.space_before = Pt(20)
p.text = "Output:"
set_run_format(p.runs[0], size=22, bold=True)
p = tf.add_paragraph()
p.text = "A single integer: Minimum initial health."
set_run_format(p.runs[0], size=18)

p = tf.add_paragraph(); p.space_before = Pt(40)
p.text = "Example Output:"
set_run_format(p.runs[0], size=22, bold=True)
p = tf.add_paragraph()
p.text = "7"
set_run_format(p.runs[0], size=48, bold=True, color=COLOR_ACCENT)

# Right Column: Visual Grid
start_x = Inches(5.5)
start_y = Inches(2.0)
cell_size = Inches(0.8)
gap = Inches(0.1)

grid_data = [
    [-2, -3, 3],
    [-5, -10, 1],
    [10, 30, -5]
]

for i in range(3):
    for j in range(3):
        val = grid_data[i][j]
        
        # Create Cell Shape
        x = start_x + (j * (cell_size + gap))
        y = start_y + (i * (cell_size + gap))
        shape = slide.shapes.add_shape(1, x, y, cell_size, cell_size)
        
        # Flat Color Logic
        if val < 0:
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_RED
            text_color = RGBColor(255,255,255)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_GREEN
            text_color = RGBColor(0,0,0)

        # Add Number
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = str(val)
        p.alignment = PP_ALIGN.CENTER
        set_run_format(p.runs[0], size=18, bold=True, color=text_color)
        
        # Center text vertically
        tf.margin_top = Inches(0.2)

# ==========================================
# SLIDE 4: Proposed Algorithm
# ==========================================
slide = prs.slides.add_slide(slide_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_BG

# Flat Header
header_shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(1.0))
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = COLOR_ACCENT
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Proposed Algorithm"
set_run_format(p.runs[0], size=24, bold=True, color=RGBColor(255,255,255))

# Left Column: Steps
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3.5), Inches(5))
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = "Dynamic Programming (Bottom-Up)"
set_run_format(p.runs[0], size=20, bold=True, color=COLOR_ACCENT)
p.space_after = Pt(10)

steps = [
    "Create DP table dp[m][n]",
    "dp[i][j] stores min health to enter cell (i, j)",
    "Start from bottom-right cell",
    "Calculate min health required from next moves",
    "Ensure health is at least 1",
    "Final answer: dp[0][0]"
]

for step in steps:
    p = tf.add_paragraph()
    p.text = step
    p.level = 0
    set_run_format(p.runs[0], size=16)
    p.space_before = Pt(5)

# Right Column: Code Block
code_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(5.0), Inches(5.0))
code_box.fill.solid()
code_box.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light gray code background
tf = code_box.text_frame
tf.word_wrap = True

code_lines = [
    "for each cell (i,j) from end:",
    "  need = min(dp[i+1][j], ",
    "            dp[i][j+1])",
    "  dp[i][j] = max(1, ",
    "             need - dungeon[i][j])"
]

for line in code_lines:
    p = tf.add_paragraph()
    p.text = line
    set_run_format(p.runs[0], size=14, font_name="Courier New")
    p.space_before = Pt(2)

# ==========================================
# SAVE
# ==========================================
prs.save('Dungeon_Game_Proposal.pptx')
print("Presentation generated successfully: Dungeon_Game_Proposal.pptx")